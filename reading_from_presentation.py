import pptx
import slide
from functools import reduce
count_of_slides=0

def read_from_slide(slide_text):
    """
     Creates a slide object from the given slide text.

     :param slide_text: The text content of a slide in the presentation.
     :returns: A slide object with fields for the title, paragraphs, and slide number.
     """
    global count_of_slides
    text_of_slide=''
    if (slide_text.shapes.title != None):
        title = slide_text.shapes.title.text
        has_title=True
    else:
        title=''
        has_title=False

    for shape in slide_text.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if has_title:
                has_title=False
                continue
            text_of_slide=reduce(lambda text, text_to_add: text + text_to_add.text + '\n' if text_to_add.text != '' else '', paragraph.runs, '')

    if title=='' and text_of_slide=='': return None
    count_of_slides = count_of_slides + 1
    return slide.Slide(title,text_of_slide,count_of_slides)




